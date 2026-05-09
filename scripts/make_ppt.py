#!/usr/bin/env python3
"""毕业设计答辩PPT生成器 — YOLO26无人机航拍检测跟踪"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ─── 颜色方案 ───
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_PURPLE = RGBColor(0x81, 0x8C, 0xF8)
ACCENT_PINK = RGBColor(0xF4, 0x72, 0xB6)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)
ACCENT_YELLOW = RGBColor(0xFB, 0xBF, 0x24)
TEXT_WHITE = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ─── 辅助函数 ───
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_gradient_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    return shape

def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=TEXT_WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(6)
        p.level = 0
    return tf

def slide_header(slide, title, subtitle=None):
    add_gradient_bar(slide, 0, 0, 13.333, 0.06)
    add_text_box(slide, 0.8, 0.3, 11.7, 0.7, title, font_size=32, bold=True, color=ACCENT_BLUE)
    if subtitle:
        add_text_box(slide, 0.8, 0.9, 11.7, 0.5, subtitle, font_size=16, color=TEXT_GRAY)
    # divider
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.01))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x33, 0x41, 0x55); shape.line.fill.background()

# ═══════════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_gradient_bar(slide, 0, 0, 13.333, 0.08)
add_text_box(slide, 1.5, 1.8, 10.3, 1.2,
    '基于改进YOLO26的无人机航拍\n多目标检测与跟踪研究',
    font_size=40, bold=True, color=WHITE)
add_text_box(slide, 1.5, 3.3, 10.3, 0.6,
    'YOLO26-P2-Tracking: ECA + CoordAtt + DySample + WeightedConcat',
    font_size=18, color=ACCENT_BLUE)
add_text_box(slide, 1.5, 4.3, 10.3, 0.5,
    'VisDrone2019-MOT 数据集 · Person + Vehicle 双类检测 · 200轮训练',
    font_size=16, color=TEXT_GRAY)

# Bottom info
add_text_box(slide, 1.5, 5.8, 5, 0.5, '毕业设计答辩', font_size=20, bold=True, color=TEXT_WHITE)
add_text_box(slide, 1.5, 6.3, 5, 0.4, '指导教师: XXX   |   答辩时间: 2026年5月', font_size=14, color=TEXT_GRAY)

# ═══════════════════════════════════════════
# Slide 2: 目录
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '目 录', 'CONTENTS')

toc = [
    ('01', '研究背景与动机', '无人机航拍检测的挑战与应用场景'),
    ('02', '数据集与预处理', 'VisDrone2019-MOT → YOLO格式 → Person+Vehicle'),
    ('03', '标准YOLO26基线', 'YOLO26n 架构分析与训练'),
    ('04', 'P2-Tracking改进模型', '六大自定义模块：ECA/CoordAtt/DySample/WeightedConcat'),
    ('05', '注意力模块详解', 'ECA · CoordAtt · DySample · BiFPN原理'),
    ('06', '训练策略与结果', '两阶段训练 · 200轮 · mAP分析'),
    ('07', '总结与展望', '消融实验计划 · Tracking头开发'),
]
for i, (num, title, desc) in enumerate(toc):
    y = 1.8 + i * 0.72
    add_text_box(slide, 1.2, y, 0.8, 0.5, num, font_size=24, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, 2.0, y, 5, 0.5, title, font_size=20, bold=True, color=TEXT_WHITE)
    add_text_box(slide, 2.0, y + 0.32, 8, 0.35, desc, font_size=13, color=TEXT_GRAY)

# ═══════════════════════════════════════════
# Slide 3: 研究背景
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '01  研究背景与动机', 'Background & Motivation')

# Problem cards
for i, (icon, title, desc) in enumerate([
    ('🎯', '目标极小', '行人占图像0.02%~0.5%\n标准模型P3/8头最小检测8px'),
    ('📐', '视角多变', '无人机俯拍视角\n目标尺度、角度剧烈变化'),
    ('⚡', '实时性要求', '机载边缘计算场景\n需在精度与速度间平衡'),
    ('🔍', '多目标跟踪', '检测是跟踪的基础\n需高质量检测结果支撑MOT'),
]):
    left = 0.8 + i * 3.05
    add_card(slide, left, 1.7, 2.85, 2.0)
    add_text_box(slide, left + 0.2, 1.85, 2.5, 0.5, f'{icon}  {title}', font_size=22, bold=True, color=ACCENT_BLUE)
    add_text_box(slide, left + 0.2, 2.4, 2.45, 1.1, desc, font_size=14, color=TEXT_GRAY)

# Solution
add_card(slide, 0.8, 4.0, 11.7, 2.8, RGBColor(0x1A, 0x2E, 0x48))
add_text_box(slide, 1.1, 4.15, 11.2, 0.5, '💡 改进方案：YOLO26-P2-Tracking', font_size=22, bold=True, color=ACCENT_GREEN)
items = [
    '• P2/4 检测头：最小检测尺寸从 8×8px 降至 4×4px，专攻无人机小目标',
    '• ECA 通道注意力：~0 额外参数，自适应调整通道重要性',
    '• CoordAtt 坐标注意力：保留精确空间位置，提升小目标定位精度',
    '• DySample 动态上采样：内容自适应插值，优于最近邻/双线性',
    '• WeightedConcat 加权融合：BiFPN风格，学习最佳多层特征融合比例',
]
add_bullet_frame(slide, 1.1, 4.7, 11.2, 1.8, items, font_size=15, color=TEXT_WHITE)

# ═══════════════════════════════════════════
# Slide 4: 数据集
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '02  数据集与预处理', 'Dataset & Preprocessing')

add_card(slide, 0.8, 1.7, 5.6, 2.4)
add_text_box(slide, 1.1, 1.85, 5, 0.5, 'VisDrone2019-MOT', font_size=22, bold=True, color=ACCENT_BLUE)
items = [
    '• 来源：Tianjin University 无人机航拍',
    '• 训练集：56个序列 / 24,198张图片',
    '• 验证集：7个序列 / 2,846张图片',
    '• 原始类别：12类 (含ignored/others)',
    '• 标注格式：MOT (frame,id,x,y,w,h,score,cat,trunc,occ)',
    '• 总量：116万+标注框',
]
add_bullet_frame(slide, 1.1, 2.3, 5, 1.5, items, font_size=14)

add_card(slide, 6.8, 1.7, 5.7, 2.4)
add_text_box(slide, 7.1, 1.85, 5, 0.5, '类别映射 → 2类检测', font_size=22, bold=True, color=ACCENT_GREEN)
add_text_box(slide, 7.1, 2.3, 5.2, 1.5, '''person (328,701个)
  ← pedestrian(1) + people(2)

vehicle (592,392个)
  ← car(4) + van(5) + truck(6) + bus(9)

丢弃 (248,266个)
  ← ignored/bicycle/tricycle/awning-tricycle/
     motor/others''', font_size=14, color=TEXT_WHITE)

add_card(slide, 0.8, 4.4, 11.7, 2.5)
add_text_box(slide, 1.1, 4.55, 11, 0.5, '数据预处理流程', font_size=20, bold=True, color=ACCENT_YELLOW)
add_text_box(slide, 1.1, 5.0, 11, 1.5, '''MOT标注(.txt/per_seq) ──→  按帧分组 ──→  PIL读取尺寸 ──→  类别映射(person/vehicle) ──→  YOLO归一化格式
                                                                  ↓
                                              丢弃(ignored, bicycle, tricycle, motor等)''',
    font_size=15, color=TEXT_WHITE)

# ═══════════════════════════════════════════
# Slide 5: 标准YOLO26基线
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '03  标准YOLO26n 基线模型', 'Baseline Model')

add_card(slide, 0.8, 1.7, 5.6, 3.0)
add_text_box(slide, 1.1, 1.85, 5, 0.5, '模型规格', font_size=22, bold=True, color=ACCENT_BLUE)
add_text_box(slide, 1.1, 2.3, 5, 2.2, '''参数量：      2,504,580 (2.50M)
GFLOPs：      5.8
层数：          260
检测头：      3层 (P3/8, P4/16, P5/32)
预训练：      yolo26n.pt (COCO)
优化器：      SGD (lr=0.01, cos调度)
Batch Size： 8 (8GB显存)
Epochs：      60 → 200 (自动衔接)
输入尺寸：  512×512''', font_size=14, color=TEXT_WHITE)

add_card(slide, 6.8, 1.7, 5.7, 3.0)
add_text_box(slide, 7.1, 1.85, 5, 0.5, '训练配置', font_size=22, bold=True, color=ACCENT_GREEN)
add_text_box(slide, 7.1, 2.3, 5.2, 2.2, '''数据集：  VisDrone2019-MOT (2类)
数据盘：  D:\\yolo26_cache\\ (SSD)
缓存：      disk cache (.npy)
增强：      mosaic, flip, hsv, scale
Warmup：  3 epochs
早停：      patience=15
AMP：       False
Workers：  4
输出：      D:\\yolo26_cache\\runs\\detect\\''', font_size=14, color=TEXT_WHITE)

add_card(slide, 0.8, 4.95, 11.7, 1.9)
add_text_box(slide, 1.1, 5.1, 11, 0.5, '当前训练进度', font_size=20, bold=True, color=ACCENT_YELLOW)
add_text_box(slide, 1.1, 5.55, 11, 1.0, '''Epoch 16 | mAP@50: 0.5120 | mAP@50-95: 0.2535 | Precision: 0.7077 | Recall: 0.4661
Phase 1: 60轮 → Phase 2: 自动衔接200轮    GPU: RTX 4060 Ti 8GB, ~5min/epoch''',
    font_size=15, color=TEXT_WHITE)

# ═══════════════════════════════════════════
# Slide 6: P2-Tracking 改进总览
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '04  P2-Tracking 改进模型总览', 'Improved Architecture Overview')

# 3-column comparison
cols = [
    ('标准 YOLO26n', '3检测头\n无注意力\n最近邻上采样\n直接Concat', False),
    ('YOLO26n-P2', '+P2/4检测头\n无注意力\n最近邻上采样\n直接Concat', False),
    ('P2-Tracking (改进)', '+P2/4检测头\nECA+CoordAtt\nDySample\nWeightedConcat', True),
]
for i, (title, desc, highlight) in enumerate(cols):
    left = 0.8 + i * 4.1
    border_color = ACCENT_GREEN if highlight else RGBColor(0x33, 0x41, 0x55)
    add_card(slide, left, 1.7, 3.9, 2.6, BG_CARD if not highlight else RGBColor(0x0F, 0x2A, 0x1A))
    add_text_box(slide, left + 0.2, 1.85, 3.5, 0.5, title, font_size=20, bold=True,
                 color=ACCENT_GREEN if highlight else TEXT_WHITE)
    add_text_box(slide, left + 0.2, 2.3, 3.5, 1.7, desc, font_size=14, color=TEXT_WHITE)

# 改进模块6宫格
modules = [
    ('ECA', '通道注意力\nECCV 2020', ACCENT_BLUE),
    ('CoordAtt', '坐标注意力\nCVPR 2021', ACCENT_PURPLE),
    ('DySample', '动态上采样\nICCV 2023', ACCENT_PINK),
    ('Weighted\nConcat', '加权融合\nCVPR 2020', ACCENT_YELLOW),
    ('ECA\nBottleneck', 'ECA残差块', ACCENT_GREEN),
    ('C3k2_ECA', 'ECA增强CSP', ACCENT_BLUE),
]
for i, (name, paper, color) in enumerate(modules):
    row, col = i // 3, i % 3
    left = 0.8 + col * 4.1
    top = 4.6 + row * 1.3
    add_card(slide, left, top, 3.9, 1.1, RGBColor(0x1A, 0x24, 0x36))
    add_text_box(slide, left + 0.2, top + 0.1, 2.0, 0.8, name, font_size=16, bold=True, color=color)
    add_text_box(slide, left + 2.2, top + 0.15, 1.5, 0.8, paper, font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════
# Slide 7: 骨干网络对比
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '04  骨干网络架构对比', 'Backbone Architecture Comparison')

# Standard backbone
add_card(slide, 0.8, 1.7, 5.6, 4.5, BG_CARD)
add_text_box(slide, 1.1, 1.85, 5, 0.5, '标准 YOLO26n Backbone', font_size=20, bold=True, color=TEXT_GRAY)
add_text_box(slide, 1.1, 2.3, 5.1, 3.5, '''P1/2  Conv(64)
P2/4  Conv(128)
       C3k2(256)
P3/8  Conv(256)
       C3k2(512)
P4/16 Conv(512)
       C3k2(512)
P5/32 Conv(1024)
       C3k2(1024)
       SPPF(1024)
       C2PSA(1024)''', font_size=14, color=TEXT_WHITE, font_name='Consolas')

# Improved backbone
add_card(slide, 6.8, 1.7, 5.7, 4.5, RGBColor(0x0F, 0x2A, 0x1A))
add_text_box(slide, 7.1, 1.85, 5, 0.5, 'P2-Tracking Backbone', font_size=20, bold=True, color=ACCENT_GREEN)
add_text_box(slide, 7.1, 2.3, 5.1, 3.5, '''P1/2  Conv(64)
P2/4  Conv(128)
       C3k2_ECA(256)   ← ECA注意力
       CoordAtt         ← 坐标注意力
P3/8  Conv(256)
       C3k2_ECA(512)   ← ECA注意力
       CoordAtt         ← 坐标注意力
P4/16 Conv(512)
       C3k2(512)
P5/32 Conv(1024)
       C3k2(1024)
       SPPF(1024)
       C2PSA(1024)''', font_size=14, color=TEXT_WHITE, font_name='Consolas')

add_card(slide, 0.8, 6.4, 11.7, 0.6, RGBColor(0x1A, 0x2E, 0x48))
add_text_box(slide, 1.1, 6.5, 11, 0.4, '🔑 关键：ECA+CoordAtt 仅部署在 P2 和 P3 高分辨率层，小目标信息最丰富，注意力收益最大。P4/P5 由 C2PSA 自注意力覆盖。',
    font_size=14, color=ACCENT_YELLOW)

# ═══════════════════════════════════════════
# Slide 8: Neck 对比
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '04  Neck (FPN+PAN) 对比', 'Feature Pyramid Neck Comparison')

# Standard neck
add_card(slide, 0.8, 1.7, 5.6, 3.8, BG_CARD)
add_text_box(slide, 1.1, 1.85, 5, 0.5, '标准 Neck', font_size=20, bold=True, color=TEXT_GRAY)
add_text_box(slide, 1.1, 2.3, 5.1, 3.0, '''上采样: nn.Upsample(nearest)
         直接复制像素 → 模糊
融合:    Concat(list)
         等权重拼接 → 盲目

3个检测头: P3, P4, P5''', font_size=14, color=TEXT_WHITE, font_name='Consolas')

# Improved neck
add_card(slide, 6.8, 1.7, 5.7, 3.8, RGBColor(0x0F, 0x2A, 0x1A))
add_text_box(slide, 7.1, 1.85, 5, 0.5, 'P2-Tracking Neck', font_size=20, bold=True, color=ACCENT_GREEN)
add_text_box(slide, 7.1, 2.3, 5.1, 3.0, '''上采样: DySample(scale=2)
         内容自适应 → 保留细节
融合:    WeightedConcat(w1,w2)
         softmax学习权重 → 智能融合

4个检测头: P2, P3, P4, P5
  P2/4  → 极小目标 (4px+) ← 核心''', font_size=14, color=TEXT_WHITE, font_name='Consolas')

# Comparison table
add_card(slide, 0.8, 5.7, 11.7, 1.3)
add_text_box(slide, 1.1, 5.85, 11, 0.5, '对比总结', font_size=18, bold=True, color=ACCENT_YELLOW)
add_text_box(slide, 1.1, 6.25, 11, 0.6, '检测尺度: 3→4 (+P2)  │  上采样: 最近邻→DySample  │  融合: Concat→WeightedConcat  │  注意力: 无→ECA+CoordAtt  │  参数: 2.50M→2.80M (+12%)',
    font_size=14, color=TEXT_WHITE)

# ═══════════════════════════════════════════
# Slide 9: ECA + CoordAtt
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '05  注意力模块详解 (1/2)', 'ECA + CoordAtt')

add_card(slide, 0.8, 1.7, 5.8, 4.8)
add_text_box(slide, 1.1, 1.85, 5, 0.5, '🔵 ECA — 高效通道注意力', font_size=22, bold=True, color=ACCENT_BLUE)
add_text_box(slide, 1.1, 2.35, 5.3, 0.3, 'ECCV 2020 · ECA-Net · conv.py:676', font_size=12, color=TEXT_GRAY)
add_text_box(slide, 1.1, 2.7, 5.3, 3.5, '''核心思想：
  用 1D 卷积替代全连接层做通道交互

自适应 Kernel Size：
  k = |log₂(C)/2 + 0.5|  (取奇数)
  通道数 64→k=3, 256→k=5, 1024→k=5

计算流程：
  AdaptiveAvgPool2d(1) → [C,1,1]
  → Conv1d(k)          → 跨通道交互
  → Sigmoid()          → 归一化权重
  → × 原特征           → 通道重标定

参数开销：仅 k 个 float (≈3~5个)
   vs SE模块的 2×C²/r 参数

部署：Backbone P2/P3 的 C3k2_ECA 中''', font_size=13, color=TEXT_WHITE)

add_card(slide, 7.0, 1.7, 5.5, 4.8)
add_text_box(slide, 7.3, 1.85, 5, 0.5, '🟣 CoordAtt — 坐标注意力', font_size=22, bold=True, color=ACCENT_PURPLE)
add_text_box(slide, 7.3, 2.35, 5, 0.3, 'CVPR 2021 · conv.py:707', font_size=12, color=TEXT_GRAY)
add_text_box(slide, 7.3, 2.7, 5, 3.5, '''核心思想：
  将 2D 池化分解为 X/Y 两个 1D 编码
  保留精确空间坐标信息

计算流程：
  X方向池化([H,1]) ⊕ Y方向池化([1,W])
  → 拼接 → Conv1×1降维(r=32)
  → BN+Hardswish → 分离X/Y
  → 双Sigmoid门控 → a_h × a_w × 原特征

优势：
  SE注意力：抹掉空间信息
  CoordAtt： 保留 (x,y) 坐标
  → 小目标定位精度显著提升

部署：Backbone P2/P3 层单独使用''', font_size=13, color=TEXT_WHITE)

# ═══════════════════════════════════════════
# Slide 10: DySample + WeightedConcat
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '05  注意力模块详解 (2/2)', 'DySample + WeightedConcat + ECABottleneck')

add_card(slide, 0.8, 1.7, 3.7, 4.8)
add_text_box(slide, 1.1, 1.85, 3.2, 0.5, '🟣 DySample', font_size=20, bold=True, color=ACCENT_PINK)
add_text_box(slide, 1.1, 2.25, 3.2, 0.3, 'ICCV 2023 · conv.py:771', font_size=11, color=TEXT_GRAY)
add_text_box(slide, 1.1, 2.6, 3.2, 3.6, '''动态上采样

Pool(8×8)→Conv→SiLU
→Conv→生成offset

offset缩放0.25/Max(W,1)
限制偏移范围

grid_sample(特征,偏移网格)
内容自适应采样

vs 最近邻：
保留更多小目标细节''', font_size=13, color=TEXT_WHITE)

add_card(slide, 4.8, 1.7, 3.7, 4.8)
add_text_box(slide, 5.1, 1.85, 3.2, 0.5, '🟡 WeightedConcat', font_size=20, bold=True, color=ACCENT_YELLOW)
add_text_box(slide, 5.1, 2.25, 3.2, 0.3, 'CVPR 2020 (BiFPN) · conv.py:745', font_size=11, color=TEXT_GRAY)
add_text_box(slide, 5.1, 2.6, 3.2, 3.6, '''加权特征融合

w = softmax(params)
learnable weights

feat₁×w₁ + feat₂×w₂
→ torch.cat

例: P3=0.7, P4=0.3
自动学习最佳比例

取代所有 Concat 层''', font_size=13, color=TEXT_WHITE)

add_card(slide, 8.8, 1.7, 3.7, 4.8)
add_text_box(slide, 9.1, 1.85, 3.2, 0.5, '🟢 ECABottleneck', font_size=20, bold=True, color=ACCENT_GREEN)
add_text_box(slide, 9.1, 2.25, 3.2, 0.3, 'block.py:486', font_size=11, color=TEXT_GRAY)
add_text_box(slide, 9.1, 2.6, 3.2, 3.6, '''ECA增强残差块

Conv(1×1)→Conv(3×3)
→ECA→+残差连接

vs 标准Bottleneck：
仅增加~5个参数

被 C3k2_ECA 调用
替代标准 Bottleneck

部署：Backbone P2/P3''', font_size=13, color=TEXT_WHITE)

# ═══════════════════════════════════════════
# Slide 11: 训练策略
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '06  训练策略与当前结果', 'Training Strategy & Results')

# Two-phase training
add_card(slide, 0.8, 1.7, 5.6, 2.2, BG_CARD)
add_text_box(slide, 1.1, 1.85, 5, 0.5, 'Phase 1: 标准YOLO26n (进行中)', font_size=20, bold=True, color=ACCENT_BLUE)
add_text_box(slide, 1.1, 2.25, 5.1, 1.4, '''Epochs: 60 → 自动衔接200
Batch=8, imgsz=512, SGD+cos调度
建立baseline，验证数据流水线
当前 Epoch 16: mAP@50=0.512''', font_size=14, color=TEXT_WHITE)

add_card(slide, 6.8, 1.7, 5.7, 2.2, RGBColor(0x0F, 0x2A, 0x1A))
add_text_box(slide, 7.1, 1.85, 5, 0.5, 'Phase 2: P2-Tracking (待启动)', font_size=20, bold=True, color=ACCENT_GREEN)
add_text_box(slide, 7.1, 2.25, 5.2, 1.4, '''Epochs: 200 (yolo26s-p2-tracking.yaml)
预训练: yolo26s.pt 权重迁移
同数据集，对比baseline
预期 mAP@50: 0.60~0.68 (+5~10%)''', font_size=14, color=TEXT_WHITE)

# Current results table
add_card(slide, 0.8, 4.15, 11.7, 2.3)
add_text_box(slide, 1.1, 4.3, 11, 0.5, '训练指标走势 (Epoch 1→16)', font_size=20, bold=True, color=ACCENT_YELLOW)
add_text_box(slide, 1.1, 4.75, 11, 1.5, '''Epoch  1:  mAP@50=0.350  mAP@50-95=0.152  P=0.435  R=0.410
Epoch  5:  mAP@50=0.498  mAP@50-95=0.235  P=0.656  R=0.469
Epoch 10:  mAP@50=0.500  mAP@50-95=0.248  P=0.697  R=0.454
Epoch 16:  mAP@50=0.512  mAP@50-95=0.254  P=0.708  R=0.466
趋势: mAP稳步上升 (+16%)，Precision已稳定~0.70，Recall缓慢恢复中。中期阶段，预计200轮mAP@50可达0.60+''',
    font_size=14, color=TEXT_WHITE, font_name='Consolas')

# ═══════════════════════════════════════════
# Slide 12: 总结
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide_header(slide, '07  总结与展望', 'Summary & Future Work')

# Summary
add_card(slide, 0.8, 1.7, 11.7, 2.5)
add_text_box(slide, 1.1, 1.85, 11, 0.5, '✅ 已完成工作', font_size=22, bold=True, color=ACCENT_GREEN)
items = [
    '• VisDrone2019-MOT → YOLO 2类(person+vehicle) 数据转换流水线 (scripts/convert_mot2yolo.py)',
    '• 标准YOLO26n 基线训练启动，60→200轮自动衔接，实时中文监控 (scripts/monitor_train.py)',
    '• P2-Tracking改进模型完整架构设计：6大自定义模块 + 4检测头 + FPN+PAN改进',
    '• 知识库体系建设：HTML网页文档系统 + 训练指标详解 + 论文参考文献',
    '• 训练管理体系：D盘SSD缓存、检查点保存、自动恢复、检测点里程碑汇报',
]
add_bullet_frame(slide, 1.1, 2.3, 11, 1.6, items, font_size=14)

add_card(slide, 0.8, 4.5, 11.7, 2.5)
add_text_box(slide, 1.1, 4.65, 11, 0.5, '📋 后续计划', font_size=22, bold=True, color=ACCENT_YELLOW)
items = [
    '• Phase 2: 标准YOLO26n 200轮训练完成后，启动P2-Tracking改进模型训练',
    '• 消融实验: 逐一去掉ECA/CoordAtt/DySample/WeightedConcat，量化各模块贡献',
    '• Tracking头开发: JDE范式 (检测+ReID嵌入) + TripletLoss + CenterLoss',
    '• 多目标跟踪: 在MOT验证集上评估 MOTA/IDF1/IDs 等跟踪指标',
    '• 论文撰写: 基于实验数据和消融研究，完成毕业论文章节',
]
add_bullet_frame(slide, 1.1, 5.05, 11, 1.6, items, font_size=14)

# ═══════════════════════════════════════════
# Slide 13: 致谢
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_gradient_bar(slide, 0, 3.0, 13.333, 0.08)
add_text_box(slide, 1.5, 2.0, 10.3, 0.8, '感谢聆听', font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.0, 10.3, 0.6, '基于改进YOLO26的无人机航拍多目标检测与跟踪研究', font_size=18, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.0, 10.3, 0.5, 'VisDrone2019-MOT · ECA · CoordAtt · DySample · WeightedConcat', font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.5, 10.3, 0.5, '欢迎提问', font_size=24, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

# ─── 保存 ───
output_path = 'docs/knowledge/graduation_defense.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
